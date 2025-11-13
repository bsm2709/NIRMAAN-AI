from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


def add_title_slide(prs: Presentation, title: str, subtitle: str = "") -> None:
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    if subtitle:
        slide.placeholders[1].text = subtitle


def add_bullet_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for i, b in enumerate(bullets):
        p = body.add_paragraph() if i > 0 else body.paragraphs[0]
        p.text = b
        p.level = 0


def add_two_column_slide(prs: Presentation, title: str, left_items: list[str], right_items: list[str]) -> None:
    # Use blank slide and draw title + two text boxes manually for compatibility
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(12.0), Inches(0.6))
    tp = title_box.text_frame.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(24)
    tp.font.bold = True

    left_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(6.0), Inches(4.5))
    right_box = slide.shapes.add_textbox(Inches(6.9), Inches(1.3), Inches(6.0), Inches(4.5))

    lt = left_box.text_frame
    rt = right_box.text_frame
    lt.clear(); rt.clear()
    for i, t in enumerate(left_items):
        p = lt.add_paragraph() if i > 0 else lt.paragraphs[0]
        p.text = f"• {t}"
        p.font.size = Pt(16)
    for i, t in enumerate(right_items):
        p = rt.add_paragraph() if i > 0 else rt.paragraphs[0]
        p.text = f"• {t}"
        p.font.size = Pt(16)


def build_presentation(output_path: str = "NirmaanAI_Synopsis.pptx") -> None:
    prs = Presentation()

    # Custom cover (with institute and team details)
    cover_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(cover_layout)
    box = slide.shapes.add_textbox(Inches(0.7), Inches(0.7), Inches(12.8), Inches(1.2))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "AI-Based Delay Prediction and Transparency System for Public Project Management (Nirmaan.ai)"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(37, 99, 235)

    sub = slide.shapes.add_textbox(Inches(0.7), Inches(1.9), Inches(12.8), Inches(0.8))
    p2 = sub.text_frame.paragraphs[0]
    p2.text = "Acharya Institute of Technology — Dept. of CSE | AY 2024–2025"
    p2.font.size = Pt(18)

    team = slide.shapes.add_textbox(Inches(0.7), Inches(2.6), Inches(12.8), Inches(3.5))
    tf3 = team.text_frame
    tf3.clear()
    lines = [
        "Project Group No: __________",
        "",
        "Team Members:",
        "1. 1AY22CS040 — BASANAGOUDA METI",
        "2. 1AY23CS402 — CHANDRASHEKARA SHIVACHARYA V M",
        "3. 1AY22CS050 — CHINMAY K G",
        "4. 1AY22CS013 — ADITYA S G",
        "",
        "Project Guide: SUPREETHA H H (Assistant Professor)",
        "Date of Submission: __________",
    ]
    for i, t in enumerate(lines):
        para = tf3.add_paragraph() if i > 0 else tf3.paragraphs[0]
        para.text = t
        para.font.size = Pt(14)

    # Agenda (based on the provided template image)
    add_bullet_slide(
        prs,
        "Agenda",
        [
            "1. Introduction",
            "2. Literature Survey",
            "3. Problem Statement and Objectives",
            "4. Proposed Methodology",
            "5. Deliverables and Impact",
            "6. Group Responsibilities",
            "7. Gantt Chart",
            "8. References",
        ],
    )

    # Introduction (context + motivation)
    add_bullet_slide(
        prs,
        "Introduction",
        [
            "Public projects in India often face delays, overruns and opacity",
            "Nirmaan.ai provides transparent monitoring with role-based access",
            "ML-driven delay prediction + real-time dashboards + citizen feedback",
        ],
    )

    # Literature Survey (concise)
    add_two_column_slide(
        prs,
        "Literature Survey (Highlights)",
        [
            "IoT + ML for risk analysis (IEEE 2020)",
            "GPS/Web dashboards (IJERT 2019)",
            "ML delay prediction (Springer 2021)",
            "Blockchain transparency (Elsevier 2022)",
            "Smart city dashboards (ACM 2018)",
        ],
        [
            "Gaps: limited transparency, no public access",
            "Missing live project integration",
            "High cost/complexity for citizens",
            "Nirmaan.ai: modular, public-facing, predictive",
        ],
    )

    # Problem + Objectives
    add_bullet_slide(
        prs,
        "Problem Statement",
        [
            "Delays, budget overruns, poor monitoring and limited transparency",
            "Lack of real-time oversight and public engagement",
        ],
    )
    add_bullet_slide(
        prs,
        "Objectives",
        [
            "Real-time tracking and geo-tagged map views",
            "ML-based delay risk prediction (≥85% accuracy on train set)",
            "Role-based dashboards (admin/official/citizen)",
            "Public transparency and feedback",
            "Centralized reporting and secure JWT auth",
        ],
    )

    # Methodology
    add_bullet_slide(
        prs,
        "Proposed Methodology",
        [
            "Project data input, geo-tagging and dashboards",
            "Python ML model (timeline, progress, budget) served via API",
            "Real-time updates; status indicators (on-track, delayed, completed)",
            "RBAC: Admin (create), Official (update), Citizen (view/flag)",
            "Visualization: map + filters; logs and trend graphs",
        ],
    )

    # Requirements (Functional / Non-Functional)
    add_two_column_slide(
        prs,
        "Requirements",
        [
            "Functional: tracking, delay engine, RBAC",
            "Public portal, secure auth (JWT), feedback",
        ],
        [
            "Non-Functional: performance, scalability",
            "Cross-browser, security (HTTPS), accuracy",
        ],
    )

    # Tech stack
    add_bullet_slide(
        prs,
        "Technology Stack",
        [
            "Frontend: React",
            "Backend: Flask API (current codebase)",
            "DB: SQLite (dev) / pluggable",
            "ML: TensorFlow/keras-based model",
            "Maps: Leaflet (current) / Mapbox (option)",
        ],
    )

    # Deliverables & Use cases
    add_bullet_slide(
        prs,
        "Deliverables & Use Cases",
        [
            "Role-based dashboards + secure login",
            "Interactive project map and logs",
            "Delay prediction API + results",
            "Citizen issue reporting",
            "Use cases: official updates; citizen verification",
        ],
    )

    # Responsibilities
    add_two_column_slide(
        prs,
        "Group Responsibilities",
        ["Data/ML: model training & API", "Backend: auth, projects, comments"],
        ["Frontend: dashboards, map, UX", "Docs & testing: reports, PPT, QA"],
    )

    # Gantt chart (textual)
    add_bullet_slide(
        prs,
        "Gantt (Planned Milestones)",
        [
            "Req. Analysis → Design → ML integration",
            "Dashboard/Map → Testing/Docs → Submission",
        ],
    )

    # References
    add_bullet_slide(
        prs,
        "References",
        [
            "IEEE 2020 (IoT+ML), IJERT 2019 (GPS/Web)",
            "Springer 2021 (ML delays), Elsevier 2022 (Blockchain)",
            "Mapbox, Joblib, JWT, Node/Express docs",
        ],
    )

    prs.save(output_path)
    print(f"Saved: {output_path}")


if __name__ == "__main__":
    build_presentation()


